import os
import hashlib
from pathlib import Path
from typing import Dict, List, Optional
from django.conf import settings
from django.core.files.storage import default_storage
from docling.document_converter import DocumentConverter
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling.document_converter import ConversionResult
import pypdf as PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

class DocumentProcessor:
    """Serviço para processamento de documentos usando Docling"""
    
    def __init__(self):
        self.converter = DocumentConverter()
        
        # Configurações otimizadas para português
        self.pdf_options = PdfPipelineOptions()
        self.pdf_options.do_ocr = True
        self.pdf_options.ocr_options.lang = ["por", "eng"]  # Português e inglês
        
    def process_document(self, document_path: str, output_dir: str) -> Dict:
        """Processa documento e converte para markdown"""
        try:
            # Determinar formato do arquivo
            file_extension = Path(document_path).suffix.lower()
            
            if file_extension == '.pdf':
                return self._process_pdf(document_path, output_dir)
            elif file_extension in ['.docx', '.doc']:
                return self._process_word(document_path, output_dir)
            elif file_extension in ['.xlsx', '.xls']:
                return self._process_excel(document_path, output_dir)
            elif file_extension in ['.pptx', '.ppt']:
                return self._process_powerpoint(document_path, output_dir)
            elif file_extension in ['.txt', '.md']:
                return self._process_text(document_path, output_dir)
            else:
                raise ValueError(f"Formato de arquivo não suportado: {file_extension}")
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown_content': '',
                'metadata': {}
            }
    
    def _process_pdf(self, document_path: str, output_dir: str) -> Dict:
        """Processa PDF usando Docling"""
        try:
            # Usar Docling para conversão
            result: ConversionResult = self.converter.convert(
                document_path,
                pipeline_options=self.pdf_options
            )
            
            # Extrair markdown
            markdown_content = result.document.export_to_markdown()
            
            # Extrair metadados
            metadata = {
                'pages': len(result.document.pages) if hasattr(result.document, 'pages') else 0,
                'title': getattr(result.document, 'title', ''),
                'author': getattr(result.document, 'author', ''),
                'creation_date': getattr(result.document, 'creation_date', ''),
                'text_length': len(markdown_content),
                'processing_method': 'docling'
            }
            
            # Salvar markdown processado
            output_path = os.path.join(output_dir, f"{Path(document_path).stem}.md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'markdown_content': markdown_content,
                'output_path': output_path,
                'metadata': metadata
            }
            
        except Exception as e:
            # Fallback para PyPDF2 se Docling falhar
            return self._process_pdf_fallback(document_path, output_dir)
    
    def _process_pdf_fallback(self, document_path: str, output_dir: str) -> Dict:
        """Fallback para processamento de PDF com PyPDF2"""
        try:
            with open(document_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                text_content = []
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"## Página {page_num + 1}\n\n{text}\n")
                
                markdown_content = "\n".join(text_content)
                
                metadata = {
                    'pages': len(pdf_reader.pages),
                    'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                    'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                    'text_length': len(markdown_content),
                    'processing_method': 'pypdf2_fallback'
                }
                
                output_path = os.path.join(output_dir, f"{Path(document_path).stem}.md")
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(markdown_content)
                
                return {
                    'success': True,
                    'markdown_content': markdown_content,
                    'output_path': output_path,
                    'metadata': metadata
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown_content': '',
                'metadata': {}
            }
    
    def _process_word(self, document_path: str, output_dir: str) -> Dict:
        """Processa documentos Word"""
        try:
            print(f"🔄 Processando arquivo Word: {document_path}")
            
            # Usar Docling primeiro
            try:
                result: ConversionResult = self.converter.convert(document_path)
                markdown_content = result.document.export_to_markdown()
                print(f"✅ Docling conversão bem-sucedida: {len(markdown_content)} caracteres")
            except Exception as docling_error:
                print(f"❌ Erro no Docling: {docling_error}")
                markdown_content = ""
            
            if not markdown_content.strip():
                # Fallback para python-docx
                print("🔄 Usando fallback python-docx")
                doc = DocxDocument(document_path)
                paragraphs = []
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        # Detectar estilos de cabeçalho
                        if paragraph.style.name.startswith('Heading'):
                            level = paragraph.style.name.replace('Heading ', '')
                            if level.isdigit():
                                markdown_text = f"{'#' * int(level)} {paragraph.text}"
                            else:
                                markdown_text = f"## {paragraph.text}"
                        else:
                            markdown_text = paragraph.text
                        
                        paragraphs.append(markdown_text)
                
                markdown_content = "\n\n".join(paragraphs)
            
            metadata = {
                'paragraphs': len(doc.paragraphs) if 'doc' in locals() else 0,
                'text_length': len(markdown_content),
                'processing_method': 'docling' if 'result' in locals() and result else 'python-docx'
            }
            
            output_path = os.path.join(output_dir, f"{Path(document_path).stem}.md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'markdown_content': markdown_content,
                'output_path': output_path,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown_content': '',
                'metadata': {}
            }
    
    def _process_excel(self, document_path: str, output_dir: str) -> Dict:
        """Processa planilhas Excel"""
        try:
            workbook = openpyxl.load_workbook(document_path, read_only=True)
            markdown_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                markdown_content.append(f"# {sheet_name}\n")
                
                # Converter para tabela markdown
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell for cell in row if cell is not None):
                        row_data = [str(cell) if cell is not None else '' for cell in row]
                        rows.append(row_data)
                
                if rows:
                    # Cabeçalho da tabela
                    header = "| " + " | ".join(rows[0]) + " |"
                    separator = "|" + "|".join([" --- " for _ in rows[0]]) + "|"
                    markdown_content.append(header)
                    markdown_content.append(separator)
                    
                    # Dados da tabela
                    for row in rows[1:]:
                        row_md = "| " + " | ".join(row) + " |"
                        markdown_content.append(row_md)
                
                markdown_content.append("\n")
            
            final_content = "\n".join(markdown_content)
            
            metadata = {
                'sheets': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames,
                'text_length': len(final_content),
                'processing_method': 'openpyxl'
            }
            
            output_path = os.path.join(output_dir, f"{Path(document_path).stem}.md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(final_content)
            
            return {
                'success': True,
                'markdown_content': final_content,
                'output_path': output_path,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown_content': '',
                'metadata': {}
            }
    
    def _process_powerpoint(self, document_path: str, output_dir: str) -> Dict:
        """Processa apresentações PowerPoint"""
        try:
            presentation = Presentation(document_path)
            markdown_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                markdown_content.append(f"# Slide {slide_num}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Detectar títulos (normalmente o primeiro texto grande)
                        if shape.text.strip() and len(shape.text.strip()) < 100:
                            markdown_content.append(f"## {shape.text}\n")
                        else:
                            markdown_content.append(f"{shape.text}\n")
                
                markdown_content.append("\n---\n")
            
            final_content = "\n".join(markdown_content)
            
            metadata = {
                'slides': len(presentation.slides),
                'text_length': len(final_content),
                'processing_method': 'python-pptx'
            }
            
            output_path = os.path.join(output_dir, f"{Path(document_path).stem}.md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(final_content)
            
            return {
                'success': True,
                'markdown_content': final_content,
                'output_path': output_path,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown_content': '',
                'metadata': {}
            }
    
    def _process_text(self, document_path: str, output_dir: str) -> Dict:
        """Processa arquivos de texto simples"""
        try:
            with open(document_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Se já for markdown, manter como está
            if document_path.endswith('.md'):
                markdown_content = content
            else:
                # Converter texto simples para markdown básico
                lines = content.split('\n')
                markdown_lines = []
                
                for line in lines:
                    if line.strip():
                        markdown_lines.append(line)
                    else:
                        markdown_lines.append('')
                
                markdown_content = '\n'.join(markdown_lines)
            
            metadata = {
                'lines': len(content.split('\n')),
                'text_length': len(markdown_content),
                'processing_method': 'text_reader'
            }
            
            output_path = os.path.join(output_dir, f"{Path(document_path).stem}.md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            
            return {
                'success': True,
                'markdown_content': markdown_content,
                'output_path': output_path,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'markdown_content': '',
                'metadata': {}
            }

def calculate_file_checksum(file_path: str) -> str:
    """Calcula checksum MD5 do arquivo"""
    hash_md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()